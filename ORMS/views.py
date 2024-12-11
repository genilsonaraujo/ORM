from django.views.generic import ListView
from django.shortcuts import render, redirect, get_object_or_404
from .models import Topic, Entry, Produto, MeuModelo, SalaUps # Saida, ItemSaida #importado do arquivo models
from .forms import TopicForm, EntryForm, ProdutoForm, MeuModeloForm, SalaUpsForm, ProdutoSearchForm #,SaidaForm, ItemSaidaForm, ProdutoSearchForm
from django.http import HttpResponseRedirect
from django.urls import reverse
from django.contrib.auth.decorators import login_required#permite so quem estiver logado ter acesso as viwes
# Create your views here.
from django.http import FileResponse
#from .utils import gerar_relatorio_powerpoint

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from django.conf import settings
import os

def index(request): #pega o as informaçoes e renderiza numa pagina html
    """Pagina principal do Projetointegradors"""
    return render(request, 'ORM/index.html')# faz a requisição

#def home(request):
  #  return render(request, 'projetointegrador/home.html')  #


@login_required#RESTRICAO DA PAGINA
def topics(request):
    """mostrar todos assuntos"""
    topics = Topic.objects.order_by('date_added') #recebendo banco de dados ordenados
    context = {'topics': topics}#chave,valor
    return render(request, 'ORM/topics.html', context) #faz requisição do servidor

@login_required#RESTRICAO DA PAGINA
def home(request):
    return render(request, 'ORM/home.html')

@login_required#RESTRICAO DA PAGINA
def topic(request, topic_id):
    """mostra um unico assunto a todas suas entradas"""
    topic = Topic.objects.get(id = topic_id)#pega a informação do banco dados
    entries = topic.entry_set.order_by('-date_added')#entradas relacionadas a topic serao ordenadas.
    context = {'topic': topic, 'entries': entries} #recebe dicionario duas chaves;topic e todas relacionadas a topic.
    return render(request, 'ORM/topic.html', context)

@login_required#RESTRICAO DA PAGINA
def new_topic(request):
    if request.method == 'POST':
        form = TopicForm(request.POST)
        if form.is_valid():
            form.save()
   
# Redireciona para a página que exibe os tópicos (substitua 'topic_view' pelo nome correto da sua visualização)
            return HttpResponseRedirect(reverse('topic'))
    else:
        form = TopicForm()

    context = {'form': form}
    return render(request, 'ORM/new_topic.html', context)

@login_required#RESTRICAO DA PAGINA
def new_entry(request, topic_id):
    """acrescenta uma nova entrada """
    topic = Topic.objects.get(id=topic_id)
    if request.method != 'POST':
        #nenhm dado submetido; cria um formulario em branco
        form = EntryForm() #variavel  recebe funcao da classe TopicForm que foi importado do nosso arquivo
    else: 
        # Dados de POST submetidos; processa os dados
        form = EntryForm(data=request.POST)
        if form.is_valid():#se os dados nao form verdadeiro nao salva
           new_entry = form.save(commit=False)#salva altomaticamente usando model no banco de dados
           new_entry.topic = topic
           new_entry.save()
           return HttpResponseRedirect(reverse('topic', args=[topic_id]))#vc
    context = {'topic':topic, 'form':form} #recebe dicionario duas chaves;topic e todas relacionadas a topic.
    return render(request, 'ORM/new_entry.html', context)

def edit_entry(request, entry_id):
    """edita uma entrada existente"""
    entry = Entry.objects.get(id=entry_id)#pegua o objeto pelo id
    topic = entry.topic

    if request.method != 'POST':
    #requesicao inicial; prenche previamente com a entrada atual.
        form = EntryForm(instance=entry)
    else: 
    # dados de POST submetidos, processa os dados
        form = EntryForm(instance=entry, data=request.POST)#substitui os dados pelo que ja estava
        if form.is_valid():
            form.save()
            return HttpResponseRedirect(reverse('topic', args=[entry_id]))#argumentos exigido por topic
    context = {'entry': entry, 'topic': topic, 'form': form}
    return render(request, 'ORM/edit_entry.html', context)




@login_required#RESTRICAO DA PAGINA
def lista_produtos(request):
    #primeira busca de produtos
    produtos = Produto.objects.all()
    #retornamos o template p/ lista de produtos
    return render(request, 'ORM/lista_produtos.html', {'produtos' : produtos})

@login_required#RESTRICAO DA PAGINA
def cria_produto(request):
    if request.method == 'POST':
        form = ProdutoForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('lista_produtos')
    else:
        form = ProdutoForm()
            
    return render(request, 'ORM/cria_produto.html', {'form' : form})
def remove_produto(request, produto_id):
    produto = get_object_or_404(Produto, pk=produto_id)
    produto.delete()
    return redirect('lista_produtos')


#def editar_produto(id_produto, novo_nome, nova_categoria, novo_codigo, novo_existente, novo_sku, novo_preco, nova_quantidade):
   # try:
       # produto = Produto.objects.get(id=id_produto)

        # Atualize os campos conforme necessário
       # produto.nome = novo_nome
       # produto.categoria = nova_categoria
        #produto.codigo = novo_codigo
        #produto.existente = novo_existente
        #produto.sku = novo_sku
        #produto.preco = novo_preco
       # produto.quantidade = nova_quantidade

        # Salve as alterações no banco de dados
       # produto.save()

        #return "Produto editado com sucesso."

    #except Produto.DoesNotExist:
        #return "Produto não encontrado."

# Exemplo de uso:
# editar_produto(1, "Novo Nome", "Nova Categoria", "Novo Código", True, "Novo SKU", 19.99, 50)
def editar_produto(request, produto_id):
    produto = get_object_or_404(Produto, id=produto_id)

    if request.method == 'POST':
        form = ProdutoForm(request.POST, request.FILES, instance=produto)
        if form.is_valid():
            form.save()
            return redirect('lista_produtos')
    else:
        form = ProdutoForm(instance=produto)

    return render(request, 'ORM/editar_produto.html', {'form': form, 'produto': produto})


def criar_item(request):
    if request.method == 'POST':
        form = MeuModeloForm(request.POST, request.FILES)
        if form.is_valid():
            item = form.save()  # Salva o item no banco de dados
            return redirect('detalhes_item', pk=item.pk)  # Redireciona para a página de detalhes
    else:
        # Adicione valores iniciais aqui
        valores_iniciais = {
            'campo1': 'valor padrão',
            'campo2': 'Outro valor padrão',
        }
        form = MeuModeloForm(initial=valores_iniciais)# Passa os valores iniciais para o formulário
    
    return render(request, 'ORM/criar_item.html', {'form': form})

def detalhes_item(request, pk):
    item = get_object_or_404(MeuModelo, pk=pk)  # Busca o item pelo ID (pk)
    return render(request, 'ORM/detalhes_item.html', {'item': item})




def criar_ups(request):
    if request.method == 'POST':
        form = SalaUpsForm(request.POST, request.FILES)
        if form.is_valid():
            item = form.save()  # Salva o item no banco de dados
            return redirect('detalhes_ups', pk=item.pk)  # Redireciona para a página de detalhes
    else:
        # Adicione valores iniciais aqui
        valores_iniciais = {
            'campo1': 'valor padrão',
            'campo2': 'Outro valor padrão',
        }
        form = SalaUpsForm(initial=valores_iniciais)# Passa os valores iniciais para o formulário
    
    return render(request, 'ORM/criar_ups.html', {'form': form})

def detalhes_ups(request, pk):
    ups = get_object_or_404(SalaUps, pk=pk)  # Busca o item pelo ID (pk)
    return render(request, 'ORM/detalhes_ups.html', {'ups': ups})

#editar ups
def editar_ups(request, pk):
    ups = get_object_or_404(SalaUps, pk=pk)  # Obtém o item pelo ID

    if request.method == 'POST':
        form = SalaUpsForm(request.POST, request.FILES, instance=ups)  # Atualiza o objeto existente
        if form.is_valid():
            form.save()  # Salva as alterações no banco de dados
            return redirect('detalhes_ups', pk=ups.pk)  # Redireciona para a página de detalhes
    else:
        form = SalaUpsForm(instance=ups)  # Preenche o formulário com os dados existentes

    return render(request, 'ORM/editar_ups.html', {'form': form, 'ups': ups})

#excluir ups
def excluir_ups(request, pk):
    ups = get_object_or_404(SalaUps, pk=pk)  # Obtém o item pelo ID

    if request.method == 'POST':  # Confirmação de exclusão
        ups.delete()  # Remove o item do banco de dados
        return redirect('ups_list')  # Redireciona para a página da lista de UPS

    return render(request, 'ORM/excluir_ups.html', {'ups': ups})


def ups_list(request):
    ups_list = SalaUps.objects.all().order_by('data_hora')  # Exemplo: Ordenar por data
    return render(request, 'ORM/ups_list.html', {'ups_list': ups_list})

def ups_detail(request, pk):
    ups = get_object_or_404(SalaUps, pk=pk)
    return render(request, 'ORM/ups_detail.html', {'ups': ups})



def gerar_relatorio_powerpoint():
    prs = Presentation()
    salas = SalaUps.objects.all()

    for sala in salas:
        # Adiciona um slide com layout em branco
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Slide vazio

        # Adiciona o título no topo
        title_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(0.3), width=Inches(9), height=Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = f"Sala: {sala.sala}"
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.size = Pt(24)

        # Adiciona informações textuais na lateral esquerda
        text_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(1.5), width=Inches(4.5), height=Inches(5))
        text_frame = text_box.text_frame
        text_frame.text = (
            f"Nome DS: {sala.nome_ds}\n"
            f"Potência DS: {sala.potencia_ds}\n"
            f"Energia ETE: {sala.energia_ete}\n"
            f"Energia Portaria: {sala.energia_portaria}\n"
            f"UPS1: {sala.ups1}, Potência: {sala.potencia_ups1}\n"
            f"UPS2: {sala.ups2}, Potência: {sala.potencia_ups2}\n"
            f"Código 1: {sala.cod1}, Energia: {sala.energia_cod1}\n"
            f"Código 2: {sala.cod2}, Energia: {sala.energia_cod2}\n"
            f"Observação: {sala.observacao}\n"
            f"Técnico: {sala.nome_tecnico}\n"
            f"Data/Hora: {sala.data_hora.strftime('%d/%m/%Y %H:%M')}"
        )
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(14)

        # Adiciona imagens na parte inferior, alinhadas horizontalmente
        imagens = [
            ("Imagem DS", sala.imagem_ds),
            ("Imagem ETE", sala.imagem_ete),
            ("Imagem Portaria", sala.imagem_portaria),
            ("UPS 1", sala.imagem_ups1),
            ("UPS 2", sala.imagem_ups2),
            ("Código 1", sala.imagem_cod1),
            ("Código 1 (Zoom)", sala.imagem_cod1z),
            ("Código 2", sala.imagem_cod2),
            ("Código 2 (Zoom)", sala.imagem_cod2z),
        ]

        img_left = Inches(0.5)  # Posição inicial da imagem
        img_top = Inches(6)  # Linha das imagens
        img_width = Inches(1.5)  # Largura padrão de cada imagem

        for nome, imagem in imagens:
            if imagem and os.path.exists(os.path.join(settings.MEDIA_ROOT, imagem.name)):
                img_path = os.path.join(settings.MEDIA_ROOT, imagem.name)
                
                # Adiciona a imagem proporcional
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)

                # Adiciona o texto abaixo da imagem
                text_box = slide.shapes.add_textbox(left=img_left, top=img_top + Inches(2.1), width=img_width, height=Inches(0.5))
                text_frame = text_box.text_frame
                paragraph = text_frame.add_paragraph()
                paragraph.text = nome
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Preto
                paragraph.alignment = 1  # Centralizar o texto

                img_left += Inches(2.5)  # Adiciona espaço entre as imagens

    # Salva o PowerPoint
    relatorio_path = os.path.join(settings.BASE_DIR, "relatorios", "relatorio_salas.ppsx")
    prs.save(relatorio_path)
    return relatorio_path

def download_relatorio(request):
    relatorio_path = gerar_relatorio_powerpoint()
    return FileResponse(open(relatorio_path, 'rb'), as_attachment=True, filename='relatorio_salas.ppsx')

#@login_required#RESTRICAO DA PAGINA
#def lista_saida(request):
    #saidas = Saida.objects.all()
    #return render(request, 'projetointegrador/lista_saida.html', {'saidas': saidas})

#def adicionar_saida(request):
    #if request.method == 'POST':
        #form = SaidaForm(request.POST)
        #if form.is_valid():
            #saida = form.save()
           # return redirect('detalhes_saida', pk=saida.pk)
    #else:
        #form = SaidaForm()
    #return render(request, 'projetointegrador/form_saida.html', {'form': form})

#def detalhes_saida(request, pk):
    #saida = Saida.objects.get(pk=pk)
    #if request.method == 'POST':
        #form = ItemSaidaForm(request.POST)
        #if form.is_valid():
            #item_saida = form.save(commit=False)
            #item_saida.saida = saida
            #item_saida.save()
            #return redirect('detalhes_saida', pk=pk)
    #else:
        #form = ItemSaidaForm()
    #return render(request, 'projetointegrador/detalhes_saida.html', {'saida': saida, 'form': form})


#@login_required
#def registrar_saida(request):
    #produtos_disponiveis = Produto.objects.filter(existente=True, quantidade__gt=0)

    #if request.method == 'POST':
        #form = SaidaForm(request.POST)
        #if form.is_valid():
            #nova_saida = form.save(commit=False)
            #nova_saida.save()

            # Processar cada item da saída
            #for produto_id, quantidade in request.POST.items():
                #if produto_id.startswith('produto_'):
                    #produto_id = produto_id.replace('produto_', '')
                    #produto = get_object_or_404(Produto, pk=produto_id)
                    #quantidade = int(quantidade)

                    #if quantidade > 0 and produto.quantidade >= quantidade:
                        # Criar o item de saída
                        #item_saida, created = ItemSaida.objects.get_or_create(
                            #saida=nova_saida,
                            #produto=produto,
                            #defaults={'quantidade': quantidade}
                       # )

                        # Calcular e salvar o preço total do item
                       # item_saida.preco_total = produto.preco * quantidade
                        #item_saida.save()

                        # Atualizar a quantidade disponível do produto
                        #produto.quantidade -= quantidade
                       # produto.save()

            #return redirect('lista_saida')  # Redirecionar para a lista de saídas após registrar a saída

   # else:
        #form = SaidaForm()

    #return render(request, 'projetointegrador/form_saida.html', {'form': form, 'produtos_disponiveis': produtos_disponiveis})


#@login_required
#def remover_saida(request, saida_id):
    # Obtenha a saída pelo ID (ou retorne um erro 404 se não existir)
    #saida = get_object_or_404(Saida, pk=saida_id)

    #if request.method == 'POST':
        # Deletar a saída e todos os seus itens associados
        #itens_saida = ItemSaida.objects.filter(saida=saida)
        #for item in itens_saida:
            # Restaurar a quantidade do produto associado ao item
            #item.produto.quantidade += item.quantidade
            #item.produto.save()

        # Deletar a saída (e todos os seus itens associados) do banco de dados
        #saida.delete()

        #return redirect('lista_saida')  # Redirecionar para a lista de saídas após remover a saída

    #return render(request, 'projetointegrador/remover_saida.html', {'saida': saida})


def busca_produtos(request): #define a funçao e a logica para buscar apenas pelo produto
    form = ProdutoSearchForm()
    produtos = Produto.objects.all()  # Busca todos os produtos por padrão

    if 'query' in request.GET:
        form = ProdutoSearchForm(request.GET)
        if form.is_valid():
            query = form.cleaned_data['query']
            produtos = produtos.filter(nome__icontains=query)  # Filtra produtos pelo nome

    return render(request, 'ORM/busca.produtos.html', {'form': form, 'produtos': produtos})
def busca_produtos(request):
    form = ProdutoSearchForm(request.GET or None)
    produtos = Produto.objects.all()

    if form.is_valid():
        if form.cleaned_data['categoria']:
            produtos = produtos.filter(categoria__icontains=form.cleaned_data['categoria'])
        if form.cleaned_data['marca']:
            produtos = produtos.filter(marca__icontains=form.cleaned_data['marca'])
        if form.cleaned_data['nome']:
            produtos = produtos.filter(nome__icontains=form.cleaned_data['nome'])
        if form.cleaned_data['modelo']:
            produtos = produtos.filter(modelo__icontains=form.cleaned_data['modelo'])
        if form.cleaned_data['codigo']:
            produtos = produtos.filter(codigo__icontains=form.cleaned_data['codigo'])
        if form.cleaned_data['sku']:
            produtos = produtos.filter(sku__icontains=form.cleaned_data['sku'])
        if form.cleaned_data['preco_min']:
            produtos = produtos.filter(preco__gte=form.cleaned_data['preco_min'])
        if form.cleaned_data['preco_max']:
            produtos = produtos.filter(preco__lte=form.cleaned_data['preco_max'])

    return render(request, 'ORM/busca.produtos.html', {'form': form, 'produtos': produtos})